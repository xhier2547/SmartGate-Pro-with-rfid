from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.dml.color import RGBColor
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Color Palette ─────────────────────────────────────────────
BG       = RGBColor(0x0D, 0x11, 0x17)   # dark bg
ACCENT   = RGBColor(0x00, 0x78, 0xD4)   # bright blue
ACCENT2  = RGBColor(0x00, 0xC8, 0x8A)   # emerald green
WARN     = RGBColor(0xF5, 0x9E, 0x0B)   # amber
RED_C    = RGBColor(0xEF, 0x44, 0x44)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0x8B, 0x94, 0x9E)
DARK2    = RGBColor(0x16, 0x1B, 0x22)   # card bg
DARK3    = RGBColor(0x21, 0x26, 0x2D)   # lighter card

FONT_TH  = 'TH Sarabun New'
FONT_EN  = 'Segoe UI'

# ── Layout helpers ─────────────────────────────────────────────
blank_layout = prs.slide_layouts[6]  # blank

def new_slide():
    return prs.slides.add_slide(blank_layout)

def rect(slide, x, y, w, h, fill=BG, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape

def txt(slide, text, x, y, w, h, size=20, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, font=FONT_TH, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def accent_bar(slide, x, y, w=0.06, h=0.55, color=ACCENT):
    rect(slide, x, y, w, h, fill=color)

def card(slide, x, y, w, h, color=DARK2):
    s = rect(slide, x, y, w, h, fill=color)
    return s

def badge(slide, text, x, y, w=1.8, h=0.38, bg=ACCENT, fc=WHITE, size=13):
    rect(slide, x, y, w, h, fill=bg)
    txt(slide, text, x, y, w, h, size=size, bold=True, color=fc,
        align=PP_ALIGN.CENTER, font=FONT_EN)

# ════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 7.5, fill=BG)
# gradient strip left
rect(s, 0, 0, 0.5, 7.5, fill=ACCENT)
rect(s, 0.5, 0, 0.08, 7.5, fill=RGBColor(0x00, 0x50, 0x9A))
# decorative circle top-right
shape = s.shapes.add_shape(9, Inches(10), Inches(-1.5), Inches(5), Inches(5))
shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x00, 0x3A, 0x6A)
shape.line.fill.background()

txt(s, "SMARTGATE PRO", 0.9, 1.4, 9, 1.2, size=56, bold=True,
    color=WHITE, font=FONT_EN)
txt(s, "ระบบบริหารจัดการการเข้า-ออกโรงเรียนอัจฉริยะ", 0.9, 2.55, 10, 0.7,
    size=26, bold=False, color=RGBColor(0xA0, 0xD4, 0xFF))
txt(s, "ด้วย RFID · AI Face Recognition · IoT · LINE Notification",
    0.9, 3.15, 11, 0.6, size=18, color=GRAY)

rect(s, 0.9, 3.85, 5.5, 0.05, fill=ACCENT)

txt(s, "นายอภิลักษณ์  จันทร์แก้วเดช  |  รหัส 6610110688",
    0.9, 4.05, 10, 0.5, size=16, color=GRAY)
txt(s, "ปีการศึกษา 2567", 0.9, 4.5, 5, 0.4, size=14, color=GRAY)

badge(s, "🏆  IoT Solution", 0.9, 5.4, 2, 0.45, bg=ACCENT, size=14)
badge(s, "🤖  AI-Powered",  3.1, 5.4, 2, 0.45, bg=ACCENT2, size=14)
badge(s, "📱  LINE Ready",  5.3, 5.4, 2, 0.45, bg=WARN, fc=BG, size=14)

# ════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM
# ════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 7.5, fill=BG)
rect(s, 0, 0, 13.33, 1.1, fill=DARK2)
txt(s, "⚠️  ปัญหาที่โรงเรียนเผชิญอยู่ทุกวัน", 0.5, 0.18, 12, 0.75,
    size=28, bold=True, color=WHITE)

problems = [
    ("⏱️", "เสียเวลา 10+ นาที/วัน", "การเช็คชื่อด้วยมือทำให้เสียเวลาเรียนสะสมกว่า 30 ชั่วโมง/ปี"),
    ("📵", "ผู้ปกครองไม่รู้ข้อมูล", "ไม่มีระบบแจ้งเตือนทำให้ผู้ปกครองกังวลเรื่องความปลอดภัยลูก"),
    ("🃏", "การฝากบัตรและโกงเข้าเรียน", "ไม่มีระบบยืนยันตัวตน นักเรียนฝากบัตรกับเพื่อนได้ง่าย"),
    ("📊", "ข้อมูลกระจัดกระจาย", "ไม่มีสถิติและ Dashboard ทำให้ผู้บริหารวิเคราะห์ข้อมูลไม่ได้"),
]
for i, (icon, title, desc) in enumerate(problems):
    cx = 0.35 + (i % 2) * 6.4
    cy = 1.4 + (i // 2) * 2.7
    card(s, cx, cy, 6.1, 2.4, color=DARK2)
    rect(s, cx, cy, 6.1, 0.07, fill=RED_C)
    txt(s, icon + "  " + title, cx+0.2, cy+0.18, 5.7, 0.55,
        size=20, bold=True, color=WHITE)
    txt(s, desc, cx+0.2, cy+0.72, 5.6, 1.2, size=16, color=GRAY)

# ════════════════════════════════════════════════════════════
# SLIDE 3 — SOLUTION OVERVIEW
# ════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 7.5, fill=BG)
rect(s, 0, 0, 13.33, 1.1, fill=DARK2)
txt(s, "💡  SmartGate Pro — วิธีแก้ปัญหาครบวงจร", 0.5, 0.18, 12, 0.75,
    size=28, bold=True, color=WHITE)

solutions = [
    ("🪪", "RFID + ESP32", "แตะบัตรเสร็จใน < 1 วินาที\nไม่ต้องเรียงคิว ไม่ต้องเซ็นชื่อ", ACCENT),
    ("🤖", "AI Face Recognition", "ตรวจใบหน้าป้องกันการฝากบัตร\nBest-of-8 Frame ทนมุมเอียงได้ดี", ACCENT2),
    ("📲", "LINE Notification", "ผู้ปกครองรับแจ้งทันทีทาง LINE\nพร้อมชื่อ เวลา และสถานะ", WARN),
    ("💡", "LED + Buzzer IoT", "ไฟสีสัญญาณที่ประตูทันที\nเขียว = ผ่าน / แดง = ไม่ผ่าน", RGBColor(0xA8, 0x55, 0xFC)),
    ("📊", "Real-time Dashboard", "ครูและผู้บริหารเห็นข้อมูลสด\nวิเคราะห์สถิติย้อนหลังได้", RGBColor(0x06, 0xB6, 0xD4)),
    ("🌐", "Remote Access", "ใช้งานได้จากทุกที่ผ่าน ngrok\nรองรับ Multi-device", RGBColor(0xF4, 0x37, 0x6D)),
]
for i, (icon, title, desc, col) in enumerate(solutions):
    cx = 0.3 + (i % 3) * 4.3
    cy = 1.3 + (i // 3) * 2.8
    card(s, cx, cy, 4.0, 2.5, color=DARK2)
    rect(s, cx, cy, 4.0, 0.07, fill=col)
    txt(s, icon, cx+0.2, cy+0.18, 0.7, 0.55, size=24)
    txt(s, title, cx+0.9, cy+0.22, 3.0, 0.5, size=17, bold=True, color=WHITE)
    txt(s, desc, cx+0.2, cy+0.78, 3.6, 1.4, size=14, color=GRAY)

# ════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE
# ════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 7.5, fill=BG)
rect(s, 0, 0, 13.33, 1.1, fill=DARK2)
txt(s, "🏗️  Architecture — การทำงานของระบบ", 0.5, 0.18, 12, 0.75,
    size=28, bold=True, color=WHITE)

# Draw flow diagram
nodes = [
    (0.4,  2.8, 2.0, 1.8, "🪪\nบัตร RFID\n+ ESP32", ACCENT),
    (3.2,  2.8, 2.0, 1.8, "☁️\nMQTT\nHiveMQ", RGBColor(0x67, 0x4E, 0xA2)),
    (6.0,  2.8, 2.0, 1.8, "⚙️\nNode.js\nServer", DARK3),
    (8.8,  1.2, 2.0, 1.8, "🖥️\nWeb\nDashboard", ACCENT2),
    (8.8,  3.3, 2.0, 1.8, "📱\nLINE\nNotification", WARN),
    (8.8,  5.3, 2.0, 1.0, "🗄️  SQLite DB", RGBColor(0x0F, 0x50, 0x2A)),
]
for nx, ny, nw, nh, label, col in nodes:
    card(s, nx, ny, nw, nh, color=col)
    txt(s, label, nx, ny, nw, nh, size=15, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)

arrows = [(2.4,3.7,3.2,3.7), (5.2,3.7,6.0,3.7),
          (8.0,2.3,8.8,2.1), (8.0,3.7,8.8,4.2), (8.0,4.5,8.8,5.8)]
for x1,y1,x2,y2 in arrows:
    line = s.shapes.add_connector(1, Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    line.line.color.rgb = ACCENT
    line.line.width = Pt(2)

# MQTT feedback arrow (ESP32 ← Server)
line2 = s.shapes.add_connector(1, Inches(6.0),Inches(4.3),Inches(2.4),Inches(4.3))
line2.line.color.rgb = ACCENT2
line2.line.width = Pt(1.5)
txt(s, "← LED/Buzzer feedback", 2.6, 4.35, 3.5, 0.35, size=11,
    color=ACCENT2, font=FONT_EN)

txt(s, "① บัตรสแกน", 0.4, 4.7, 2.0, 0.3, size=12, color=GRAY, align=PP_ALIGN.CENTER)
txt(s, "② MQTT publish", 3.0, 4.7, 2.5, 0.3, size=12, color=GRAY, align=PP_ALIGN.CENTER)
txt(s, "③ ประมวลผล AI", 6.0, 4.7, 2.5, 0.3, size=12, color=GRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 5 — HARDWARE
# ════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 7.5, fill=BG)
rect(s, 0, 0, 13.33, 1.1, fill=DARK2)
txt(s, "🔧  Hardware — อุปกรณ์และการเชื่อมต่อ", 0.5, 0.18, 12, 0.75,
    size=28, bold=True, color=WHITE)

hw = [
    ("ESP32 DevKit V1", "Dual-core 240MHz, WiFi/BT", "🧠", ACCENT),
    ("ID-20LA RFID Reader", "125kHz, UART Serial, Long Range", "📡", ACCENT2),
    ("LED 3 ดวง", "GPIO 25=เขียว, 26=เหลือง, 32=แดง", "💡", WARN),
    ("Passive Buzzer", "GPIO 13, เสียงแจ้งเตือนตามผล AI", "🔊", RGBColor(0xA8, 0x55, 0xFC)),
    ("USB Webcam HD", "1080p, AI Face Recognition", "📷", RGBColor(0xF4, 0x37, 0x6D)),
    ("Windows Server", "Node.js + SQLite + face-api.js", "🖥️", RGBColor(0x06, 0xB6, 0xD4)),
]
for i, (name, spec, icon, col) in enumerate(hw):
    cx = 0.35 + (i % 3) * 4.3
    cy = 1.3 + (i // 3) * 2.8
    card(s, cx, cy, 4.0, 2.5, color=DARK2)
    rect(s, cx, cy, 0.08, 2.5, fill=col)
    txt(s, icon, cx+0.25, cy+0.3, 0.7, 0.8, size=30)
    txt(s, name, cx+1.0, cy+0.3, 2.9, 0.55, size=17, bold=True, color=WHITE)
    txt(s, spec, cx+1.0, cy+0.85, 2.8, 1.3, size=14, color=GRAY)

txt(s, "💡 ต่อขา: ESP32 GPIO16(RX) ← TX ของ ID-20LA | LED ต่อผ่าน 220Ω | Buzzer ต่อ GND",
    0.5, 6.85, 12.5, 0.4, size=13, color=GRAY, font=FONT_TH)

# ════════════════════════════════════════════════════════════
# SLIDE 6 — AI FACE RECOGNITION
# ════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 7.5, fill=BG)
rect(s, 0, 0, 13.33, 1.1, fill=DARK2)
txt(s, "🤖  AI Face Recognition — ป้องกันการฝากบัตร", 0.5, 0.18, 12, 0.75,
    size=28, bold=True, color=WHITE)

# Left: Steps
steps = [
    ("1", "บัตรสแกน → ส่ง MQTT", ACCENT),
    ("2", "โหลด Face Descriptor จาก IndexedDB", ACCENT),
    ("3", "ถ่ายภาพจากกล้อง 8 frames / 200ms", ACCENT),
    ("4", "คำนวณ Euclidean Distance ต่อ frame", ACCENT),
    ("5", "Best Distance < 0.92 → ผ่าน ✅", ACCENT2),
    ("6", "ส่ง Feedback → LED + LINE", ACCENT2),
]
for i, (num, step, col) in enumerate(steps):
    cy = 1.35 + i * 0.87
    rect(s, 0.4, cy, 0.55, 0.55, fill=col)
    txt(s, num, 0.4, cy, 0.55, 0.55, size=20, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, step, 1.1, cy+0.04, 5.5, 0.5, size=16, color=WHITE)

# Right: Stats
card(s, 7.2, 1.3, 5.7, 5.8, color=DARK2)
txt(s, "📊 ผลการทดสอบ AI", 7.4, 1.5, 5.3, 0.5, size=18,
    bold=True, color=WHITE)
stats = [
    ("ใบหน้าตรง, แสงดี", "0.28–0.38", "✅ ผ่าน", ACCENT2),
    ("หน้าเอียง ~30°", "0.42–0.58", "✅ ผ่าน", ACCENT2),
    ("แสงน้อย/กลางคืน", "0.55–0.72", "✅ ผ่าน", ACCENT2),
    ("เอียง 45°+แสงน้อย", "0.75–0.89", "✅ ผ่าน", ACCENT2),
    ("คนละคน (Imposter)", "0.93–1.00", "🚫 ปฏิเสธ", RED_C),
]
for i, (cond, dist, result, col) in enumerate(stats):
    cy = 2.1 + i * 0.95
    rect(s, 7.3, cy, 5.5, 0.75, fill=DARK3)
    txt(s, cond, 7.5, cy+0.1, 2.8, 0.5, size=13, color=GRAY)
    txt(s, dist, 10.3, cy+0.1, 1.2, 0.5, size=13, color=WHITE, font=FONT_EN)
    txt(s, result, 11.5, cy+0.1, 1.2, 0.5, size=13, bold=True, color=col)

txt(s, f"Threshold: 0.92  |  TinyFaceDetector  |  8 Frames Sampling",
    7.4, 6.85, 5.3, 0.35, size=12, color=GRAY, font=FONT_EN)

# ════════════════════════════════════════════════════════════
# SLIDE 7 — LED BUZZER 2-WAY IoT
# ════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 7.5, fill=BG)
rect(s, 0, 0, 13.33, 1.1, fill=DARK2)
txt(s, "💡  LED + Buzzer — 2-Way IoT Feedback", 0.5, 0.18, 12, 0.75,
    size=28, bold=True, color=WHITE)

feedbacks = [
    ("🟢", "ผ่าน AI (ปกติ)", "LED เขียว 1 วิ + Beep 1800Hz", "'pass'", ACCENT2),
    ("🟡", "มาสาย", "LED เหลือง 1.5 วิ + Beep 2 ครั้ง", "'late'", WARN),
    ("🔴", "AI ไม่ผ่าน", "LED แดง 2 วิ + Buzz ยาว 500Hz", "'fail'", RED_C),
    ("🔴", "บัตรไม่รู้จัก", "LED แดงกะพริบ 3× + Beep 3 ครั้ง", "'unknown'", RED_C),
    ("🟢", "เช็คเอาท์ (ออก)", "LED เขียวสั้น 0.6 วิ + Beep เบา", "'checkout'", ACCENT2),
]
for i, (led, status, action, msg, col) in enumerate(feedbacks):
    cy = 1.3 + i * 1.1
    card(s, 0.4, cy, 12.4, 0.95, color=DARK2)
    rect(s, 0.4, cy, 0.08, 0.95, fill=col)
    txt(s, led, 0.6, cy+0.12, 0.7, 0.6, size=26)
    txt(s, status, 1.45, cy+0.18, 2.5, 0.5, size=17, bold=True, color=WHITE)
    txt(s, action, 4.1, cy+0.18, 5.2, 0.5, size=15, color=GRAY)
    txt(s, msg, 9.5, cy+0.18, 3.0, 0.5, size=14, bold=True,
        color=col, font=FONT_EN, align=PP_ALIGN.CENTER)

txt(s, "MQTT Topic: smartgate/feedback  ←  Server publish → ESP32 subscribe",
    0.5, 6.85, 12.3, 0.35, size=12, color=GRAY, font=FONT_EN)

# ════════════════════════════════════════════════════════════
# SLIDE 8 — RESULTS
# ════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 7.5, fill=BG)
rect(s, 0, 0, 13.33, 1.1, fill=DARK2)
txt(s, "✅  ผลลัพธ์และตัวเลขจริง", 0.5, 0.18, 12, 0.75,
    size=28, bold=True, color=WHITE)

kpis = [
    ("< 3 วิ", "เวลาเช็คอินทั้งระบบ", ACCENT),
    ("0.92", "AI Threshold (ยืดหยุ่น)", ACCENT2),
    ("8 Frames", "Sampling ต่อการสแกน", WARN),
    ("100%", "อัตโนมัติ ไม่ต้องเรียกชื่อ", RGBColor(0xA8, 0x55, 0xFC)),
]
for i, (val, label, col) in enumerate(kpis):
    cx = 0.4 + i * 3.2
    card(s, cx, 1.3, 2.9, 2.2, color=DARK2)
    rect(s, cx, 1.3, 2.9, 0.07, fill=col)
    txt(s, val, cx, 1.6, 2.9, 1.0, size=36, bold=True,
        color=col, align=PP_ALIGN.CENTER, font=FONT_EN)
    txt(s, label, cx, 2.6, 2.9, 0.65, size=14,
        color=GRAY, align=PP_ALIGN.CENTER)

features = [
    "✅  เช็คชื่ออัตโนมัติ 100% — ไม่ต้องเรียกชื่อ ไม่ต้องเซ็น",
    "✅  AI ป้องกันฝากบัตร — ตรวจออก หน้าเอียงสูงสุด 45°",
    "✅  LINE แจ้งผู้ปกครองทันที — ชื่อ เวลา สถานะ ครบ",
    "✅  LED/Buzzer ที่ประตู — เห็นผลทันทีไม่ต้องดูหน้าจอ",
    "✅  Remote Access ผ่าน ngrok — ครูใช้งานได้จากทุกที่",
    "✅  ทำงาน ESP32 กับ Server ต่าง WiFi — ผ่าน HiveMQ MQTT",
]
for i, feat in enumerate(features):
    cy = 3.75 + (i // 2) * 0.9
    cx = 0.4 if i % 2 == 0 else 6.8
    txt(s, feat, cx, cy, 6.2, 0.6, size=15, color=WHITE)

# ════════════════════════════════════════════════════════════
# SLIDE 9 — PRICING / PROPOSAL
# ════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 7.5, fill=BG)
rect(s, 0, 0, 13.33, 1.1, fill=DARK2)
txt(s, "💼  แพ็กเกจสำหรับโรงเรียน", 0.5, 0.18, 12, 0.75,
    size=28, bold=True, color=WHITE)

packages = [
    ("🥉 Starter", "ฟรี (Self-hosted)", DARK3,
     ["ESP32 + ID-20LA 1 จุด", "Web Dashboard", "LINE Notification", "รองรับนักเรียน ≤ 500 คน"]),
    ("🥇 Pro", "ติดต่อสอบถาม", ACCENT,
     ["2 จุดสแกน (2 ประตู)", "LED/Buzzer feedback", "AI Face Recognition", "Remote Access + Support"]),
    ("🏆 Enterprise", "กำหนดตามโรงเรียน", DARK2,
     ["ไม่จำกัดจุดสแกน", "IP Camera ติดตรึง", "Server-side AI (Python)", "Mobile App ผู้ปกครอง"]),
]
for i, (name, price, col, feats) in enumerate(packages):
    cx = 0.4 + i * 4.3
    card(s, cx, 1.3, 4.0, 5.8, color=col)
    if col == ACCENT:
        rect(s, cx, 1.3, 4.0, 0.07, fill=WHITE)
        badge(s, "⭐ แนะนำ", cx+1.2, 1.22, 1.6, 0.35, bg=WHITE,
              fc=ACCENT, size=12)
    else:
        rect(s, cx, 1.3, 4.0, 0.07, fill=GRAY)
    txt(s, name, cx, 1.5, 4.0, 0.6, size=22, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, price, cx, 2.1, 4.0, 0.55, size=18, bold=True,
        color=WARN if col==ACCENT else GRAY, align=PP_ALIGN.CENTER)
    for j, feat in enumerate(feats):
        txt(s, "✓  " + feat, cx+0.2, 2.85 + j*0.75, 3.6, 0.55,
            size=14, color=WHITE if col==ACCENT else GRAY)

txt(s, "* ราคาไม่รวมอุปกรณ์ Hardware  |  ต้นทุน Hardware ต่อ 1 จุด ≈ 800–1,200 บาท",
    0.5, 7.05, 12.3, 0.35, size=12, color=GRAY)

# ════════════════════════════════════════════════════════════
# SLIDE 10 — CLOSING / CTA
# ════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 7.5, fill=BG)
rect(s, 0, 0, 0.5,  7.5, fill=ACCENT2)
rect(s, 0.5, 0, 0.08, 7.5, fill=RGBColor(0x00, 0x8A, 0x60))
shape = s.shapes.add_shape(9, Inches(9), Inches(2.5), Inches(6), Inches(6))
shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x00, 0x35, 0x25)
shape.line.fill.background()

txt(s, "พร้อมเปลี่ยนโรงเรียนของคุณ", 0.9, 1.5, 11, 1.0,
    size=38, bold=True, color=WHITE)
txt(s, "ให้เป็น Smart School แล้วหรือยัง?", 0.9, 2.4, 11, 0.9,
    size=38, bold=True, color=ACCENT2)
rect(s, 0.9, 3.4, 6, 0.06, fill=ACCENT2)

txt(s, "SmartGate Pro ช่วยประหยัดเวลาครูและให้ผู้ปกครองสบายใจ\nด้วยการแจ้งเตือน LINE แบบ Real-time ทุกครั้งที่เด็กเข้าโรงเรียน",
    0.9, 3.6, 10.5, 1.2, size=18, color=GRAY)

badge(s, "📧  smartgate@school.th", 0.9, 5.1, 3.5, 0.55, bg=DARK2, size=15)
badge(s, "🌐  Demo: localhost:3000",  4.7, 5.1, 3.5, 0.55, bg=DARK2, size=15)

txt(s, "นายอภิลักษณ์  จันทร์แก้วเดช  |  6610110688  |  ปีการศึกษา 2567",
    0.9, 6.3, 10, 0.45, size=14, color=GRAY)

# ── Save ──────────────────────────────────────────────────────
out = r"c:\Users\wator\Documents\IOTPROJECT\smart-school-api\SmartGatePro_Pitch_6610110688.pptx"
prs.save(out)
print(f"✅ บันทึกไฟล์สำเร็จ: {out}")
